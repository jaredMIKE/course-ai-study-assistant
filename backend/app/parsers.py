from __future__ import annotations

from io import BytesIO
from pathlib import Path
import csv

from fastapi import HTTPException

from .config import SUPPORTED_EXTENSIONS
from .nlp import normalize


def _decode_text(raw: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "gb18030", "gbk", "big5", "latin-1"):
        try:
            return raw.decode(enc)
        except UnicodeDecodeError:
            continue
    return raw.decode("utf-8", errors="ignore")


def _extract_pdf(raw: bytes) -> str:
    try:
        from pypdf import PdfReader
    except Exception as exc:
        raise HTTPException(status_code=500, detail="缺少 pypdf 依赖，请先 pip install -r requirements.txt") from exc
    reader = PdfReader(BytesIO(raw))
    pages = []
    for index, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"【第{index}页】\n{text}")
    return "\n\n".join(pages)


def _extract_docx(raw: bytes) -> str:
    try:
        from docx import Document
    except Exception as exc:
        raise HTTPException(status_code=500, detail="缺少 python-docx 依赖，请先 pip install -r requirements.txt") from exc
    doc = Document(BytesIO(raw))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text.strip())
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _shape_text(shape) -> list[str]:
    parts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        text = shape.text.strip()
        if text:
            parts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))
    return parts


def _extract_pptx(raw: bytes) -> str:
    try:
        from pptx import Presentation
    except Exception as exc:
        raise HTTPException(status_code=500, detail="缺少 python-pptx 依赖，请先 pip install -r requirements.txt") from exc
    prs = Presentation(BytesIO(raw))
    slides: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            parts.extend(_shape_text(shape))
        if parts:
            slides.append(f"【第{idx}页PPT】\n" + "\n".join(parts))
    return "\n\n".join(slides)


def extract_text_from_file(filename: str, raw: bytes) -> tuple[str, str]:
    safe_name = Path(filename or "上传资料").name
    ext = Path(safe_name).suffix.lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise HTTPException(
            status_code=400,
            detail=f"暂不支持 {ext or '无后缀'} 文件。支持：{', '.join(SUPPORTED_EXTENSIONS)}",
        )
    if ext in {".txt", ".md"}:
        text = _decode_text(raw)
    elif ext == ".csv":
        decoded = _decode_text(raw)
        try:
            rows = csv.reader(decoded.splitlines())
            text = "\n".join(" | ".join(row) for row in rows)
        except Exception:
            text = decoded
    elif ext == ".pdf":
        text = _extract_pdf(raw)
    elif ext == ".docx":
        text = _extract_docx(raw)
    elif ext == ".pptx":
        text = _extract_pptx(raw)
    else:
        text = _decode_text(raw)

    text = normalize(text)
    if len(text) < 20:
        raise HTTPException(
            status_code=400,
            detail="文件文本提取结果太少，可能是扫描版 PDF、图片型课件或文件内容为空。",
        )
    return safe_name, text
